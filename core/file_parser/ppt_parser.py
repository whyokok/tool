"""
Microsoft PowerPoint file parser using python-pptx.
"""
from pathlib import Path

from pptx import Presentation

from .base import BaseParser, ParseResult


class PptParser(BaseParser):
    """Parser for PowerPoint files (.pptx)."""

    @property
    def file_type(self) -> str:
        return "ppt"

    def parse(self, file_path: Path) -> ParseResult:
        try:
            prs = Presentation(file_path)
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                if slide_text:
                    text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))

            content = "\n\n".join(text_parts)
            return self._create_result(file_path, content.strip())

        except Exception as e:
            return self._create_result(file_path, "", str(e))
